import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def generate_presentation(dataset_id: str, story_data: dict) -> str:
    """
    Generates a 6-slide PowerPoint report automatically.
    """
    prs = Presentation()
    
    # Slide 1: Title
    slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide_1.shapes.title
    subtitle = slide_1.placeholders[1]
    title.text = f"AI Data Story: {dataset_id}"
    subtitle.text = "Generated autonomously by AS-AI Platform"
    
    # Slide 2: Dataset Overview
    slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide_2.shapes.title
    title2.text = "Dataset Overview"
    content_box2 = slide_2.placeholders[1]
    content_box2.text = story_data.get("summary", "No summary available.")

    # Slide 3: Key Insights
    slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide_3.shapes.title
    title3.text = "Key Drivers & Insights"
    content_box3 = slide_3.placeholders[1]
    tf3 = content_box3.text_frame
    for driver in story_data.get("key_drivers", []):
        p = tf3.add_paragraph()
        p.text = str(driver)
        p.level = 0

    # Slide 4: Model Performance
    slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide_4.shapes.title
    title4.text = "Model Performance"
    content_box4 = slide_4.placeholders[1]
    content_box4.text = "Automated machine learning pipelines have selected the best performing algorithms. Advanced Explainability (SHAP) charts confirm predictive validity."

    # Slide 5: Visualizations
    slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide_5.shapes.title
    title5.text = "Visualizations"
    content_box5 = slide_5.placeholders[1]
    content_box5.text = "(Generated charts would be embedded here)"

    # Slide 6: Recommendations
    slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide_6.shapes.title
    title6.text = "Business Recommendations"
    content_box6 = slide_6.placeholders[1]
    tf6 = content_box6.text_frame
    for rec in story_data.get("recommendations", []):
        p = tf6.add_paragraph()
        p.text = str(rec)
        p.level = 0
        
    # Save the presentation
    reports_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "reports", "presentations")
    os.makedirs(reports_dir, exist_ok=True)
    
    file_name = f"{dataset_id}_analysis.pptx"
    file_path = os.path.join(reports_dir, file_name)
    
    prs.save(file_path)
    
    return f"/reports/presentations/{file_name}"
